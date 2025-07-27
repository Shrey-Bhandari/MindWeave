from pptx import Presentation

def extract_pptx_content(filepath):
    prs = Presentation(filepath)
    slide_data = []

    for slide in prs.slides:
        title = ""
        bullets = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()

                if not title:
                    title = text
                else:
                    bullets.append(text)

        slide_data.append({
            "slide_title": title,
            "bullets": bullets
        })

    return slide_data
