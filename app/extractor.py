# app/extractor.py

from pptx import Presentation

def extract_pptx_content(filepath):
    prs = Presentation(filepath)
    slide_data = []

    for slide in prs.slides:
        title = ""
        bullets = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for i, para in enumerate(shape.text_frame.paragraphs):
                text = para.text.strip()
                if not text:
                    continue

                if not title:
                    title = text
                else:
                    bullets.append(text)

        slide_data.append({
            "slide_title": title.strip(),
            "bullets": bullets
        })

    return slide_data
